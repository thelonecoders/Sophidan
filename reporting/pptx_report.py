"""PowerPoint (.pptx) report generator built on ``python-pptx``.

Produces 16:9 decks with a subtle blue accent theme.  All slide types
specified by the project contract are supported: title, section divider,
bullet list (with multi-level tuples), native bar/line/pie chart, image,
table, and two-column comparison slides.
"""
#
# MIT License — Academic Research Suite
# Copyright (c) 2026 — see /LICENSE for full text.
#
from __future__ import annotations

import logging
import os
from typing import Any, Dict, List, Optional, Sequence, Tuple, Union

logger = logging.getLogger(__name__)


# 16:9 layout dimensions in EMU (English Metric Units).
# python-pptx uses 914400 EMU per inch; default 16:9 slide = 13.333 x 7.5 in.
_SLIDE_W_IN = 13.333
_SLIDE_H_IN = 7.5
_PRIMARY_BLUE = "2E5C8A"      # RGB hex, no '#'
_LIGHT_BLUE = "4A90D9"
_BG_TINT = "F2F6FA"
_TEXT_DARK = "222222"
_TEXT_GREY = "555555"


class PPTXReport:
    """Build a 16:9 PowerPoint deck with a subtle blue-accent theme."""

    def __init__(
        self,
        title: str,
        subtitle: Optional[str] = None,
    ) -> None:
        """Initialise the deck builder.

        Args:
            title: Deck title (used in metadata + title slide).
            subtitle: Optional subtitle shown on the title slide.
        """
        try:
            from pptx import Presentation  # noqa: WPS433
            from pptx.util import Inches, Pt, Emu  # noqa: WPS433
        except ImportError as exc:  # pragma: no cover
            logger.error("python-pptx not installed: %s", exc)
            raise
        self.title = title
        self.subtitle = subtitle
        self._prs = Presentation()
        # Force 16:9.
        self._prs.slide_width = Inches(_SLIDE_W_IN)
        self._prs.slide_height = Inches(_SLIDE_H_IN)
        self._blank_layout = self._prs.slide_layouts[6]  # blank layout.
        # Set core properties.
        try:
            cp = self._prs.core_properties
            cp.title = title
            if subtitle:
                cp.subject = subtitle
        except Exception as exc:  # pragma: no cover
            logger.debug("core properties set failed: %s", exc)
        self._Inches = Inches
        self._Pt = Pt
        self._Emu = Emu

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------
    def _add_textbox(
        self,
        slide,
        left_in: float,
        top_in: float,
        width_in: float,
        height_in: float,
        text: str = "",
        font_size: int = 18,
        bold: bool = False,
        color_hex: str = _TEXT_DARK,
        align: str = "left",
        font_name: Optional[str] = None,
    ):
        """Add a single-line text box and return the shape."""
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR  # noqa: WPS433
        tb = slide.shapes.add_textbox(
            self._Inches(left_in), self._Inches(top_in),
            self._Inches(width_in), self._Inches(height_in),
        )
        tf = tb.text_frame
        tf.word_wrap = True
        align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
        tf.vertical_anchor = MSO_ANCHOR.TOP
        try:
            tf.paragraphs[0].alignment = align_map.get(align, PP_ALIGN.LEFT)
        except Exception:  # pragma: no cover
            pass
        if text:
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text
            run.font.size = self._Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = self._hex_color(color_hex)
            if font_name:
                run.font.name = font_name
        return tb

    @staticmethod
    def _hex_color(hex_str: str):
        from pptx.dml.color import RGBColor  # noqa: WPS433
        h = hex_str.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    def _add_accent_bar(self, slide, top_in: float = 1.7) -> None:
        """Add a thin horizontal blue accent bar."""
        from pptx.enum.shapes import MSO_SHAPE  # noqa: WPS433
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            self._Inches(0.6), self._Inches(top_in),
            self._Inches(2.5), self._Inches(0.08),
        )
        try:
            bar.fill.solid()
            bar.fill.fore_color.rgb = self._hex_color(_PRIMARY_BLUE)
            bar.line.fill.background()
            bar.shadow.inherit = False
        except Exception:  # pragma: no cover
            pass

    def _new_blank_slide(self):
        """Add a blank slide and return it."""
        return self._prs.slides.add_slide(self._blank_layout)

    # ------------------------------------------------------------------
    # Slide types
    # ------------------------------------------------------------------
    def add_title_slide(
        self,
        title: str,
        subtitle: Optional[str] = None,
        author: Optional[str] = None,
        date: Optional[str] = None,
    ) -> None:
        """Add the deck's title slide.

        Args:
            title: Large title text.
            subtitle: Optional subtitle text.
            author: Optional author / presenter line.
            date: Optional date string.
        """
        slide = self._new_blank_slide()
        # Top accent band.
        from pptx.enum.shapes import MSO_SHAPE  # noqa: WPS433
        band = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            self._Inches(0), self._Inches(0),
            self._prs.slide_width, self._Inches(2.4),
        )
        try:
            band.fill.solid()
            band.fill.fore_color.rgb = self._hex_color(_PRIMARY_BLUE)
            band.line.fill.background()
            band.shadow.inherit = False
        except Exception:  # pragma: no cover
            pass
        # Title text on the band.
        self._add_textbox(
            slide, 0.6, 0.7, _SLIDE_W_IN - 1.2, 1.2,
            text=title, font_size=36, bold=True, color_hex="FFFFFF",
            align="left",
        )
        if subtitle:
            self._add_textbox(
                slide, 0.6, 1.7, _SLIDE_W_IN - 1.2, 0.6,
                text=subtitle, font_size=18, color_hex="E5EEF7",
                align="left",
            )
        # Author + date below.
        meta_lines: List[str] = []
        if author:
            meta_lines.append(author)
        if date:
            meta_lines.append(date)
        if meta_lines:
            self._add_textbox(
                slide, 0.6, 3.0, _SLIDE_W_IN - 1.2, 1.0,
                text="  |  ".join(meta_lines), font_size=16, color_hex=_TEXT_GREY,
                align="left",
            )

    def add_section_slide(self, title: str) -> None:
        """Add a section-divider slide (large title, blue accent bar)."""
        slide = self._new_blank_slide()
        # Light tinted background rectangle.
        from pptx.enum.shapes import MSO_SHAPE  # noqa: WPS433
        try:
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                self._Inches(0), self._Inches(0),
                self._prs.slide_width, self._prs.slide_height,
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = self._hex_color(_BG_TINT)
            bg.line.fill.background()
            bg.shadow.inherit = False
        except Exception:  # pragma: no cover
            pass
        self._add_accent_bar(slide, top_in=3.0)
        self._add_textbox(
            slide, 0.6, 3.2, _SLIDE_W_IN - 1.2, 1.5,
            text=title, font_size=40, bold=True, color_hex=_PRIMARY_BLUE,
            align="left",
        )

    def add_bullet_slide(
        self,
        title: str,
        bullets: Sequence[Union[str, Tuple[str, int]]],
    ) -> None:
        """Add a bullet slide with optional nested levels.

        Args:
            title: Slide title.
            bullets: List where each item is either a ``str`` (level 0 bullet)
                or a ``tuple`` of ``(text, level)``.  ``level=0`` is the top
                bullet; higher values indent further.
        """
        slide = self._new_blank_slide()
        # Title bar.
        self._add_textbox(
            slide, 0.6, 0.4, _SLIDE_W_IN - 1.2, 0.9,
            text=title, font_size=28, bold=True, color_hex=_PRIMARY_BLUE,
            align="left",
        )
        self._add_accent_bar(slide, top_in=1.25)
        # Body textbox.
        tb = slide.shapes.add_textbox(
            self._Inches(0.8), self._Inches(1.7),
            self._Inches(_SLIDE_W_IN - 1.6), self._Inches(_SLIDE_H_IN - 2.2),
        )
        tf = tb.text_frame
        tf.word_wrap = True
        first = True
        for item in bullets:
            if isinstance(item, tuple):
                text, level = item[0], int(item[1])
            else:
                text, level = str(item), 0
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.level = max(0, level)
            bullet_char = "•" if level == 0 else "–"
            run = p.add_run()
            run.text = f"{bullet_char}  {text}"
            run.font.size = self._Pt(max(14, 22 - level * 4))
            run.font.color.rgb = self._hex_color(_TEXT_DARK)
            try:
                p.space_after = self._Pt(8)
            except Exception:  # pragma: no cover
                pass

    def add_chart_slide(
        self,
        title: str,
        chart_type: str,
        data: Union[Dict[str, Sequence[Any]], Sequence[Tuple[str, Sequence[Any]]]],
        x_label: Optional[str] = None,
        y_label: Optional[str] = None,
    ) -> None:
        """Add a slide with a native PowerPoint chart.

        Args:
            title: Slide title.
            chart_type: ``'bar'``, ``'line'``, or ``'pie'``.
            data: Either a dict ``{series_name: [values...]}`` with the first
                series treated as categories, OR a list of ``(series_name,
                values)`` tuples.  For ``'pie'``, the first series is treated
                as categories and the second as values.
            x_label: Optional x-axis label (ignored for pie).
            y_label: Optional y-axis label (ignored for pie).
        """
        from pptx.chart.data import CategoryChartData  # noqa: WPS433
        from pptx.enum.chart import XL_CHART_TYPE  # noqa: WPS433

        chart_type = chart_type.lower()
        ct_map = {
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE,
            "pie": XL_CHART_TYPE.PIE,
        }
        if chart_type not in ct_map:
            logger.warning("Unknown chart_type %r; falling back to 'bar'", chart_type)
            chart_type = "bar"
        xl_type = ct_map[chart_type]

        # Normalise data into (categories, series_list).
        categories, series_list = self._normalise_chart_data(chart_type, data)
        cd = CategoryChartData()
        cd.categories = categories
        for name, values in series_list:
            cd.add_series(name, values)

        slide = self._new_blank_slide()
        self._add_textbox(
            slide, 0.6, 0.4, _SLIDE_W_IN - 1.2, 0.9,
            text=title, font_size=28, bold=True, color_hex=_PRIMARY_BLUE,
            align="left",
        )
        self._add_accent_bar(slide, top_in=1.25)
        chart_shape = slide.shapes.add_chart(
            xl_type,
            self._Inches(0.8), self._Inches(1.7),
            self._Inches(_SLIDE_W_IN - 1.6), self._Inches(_SLIDE_H_IN - 2.2),
            cd,
        )
        chart = chart_shape.chart
        try:
            chart.has_legend = chart_type == "pie"
            if chart_type == "pie":
                chart.legend.position = chart.legend.position  # default
                chart.legend.include_in_layout = False
            if chart_type != "pie" and (x_label or y_label):
                chart.has_title = False
                # Axis titles.
                if x_label:
                    try:
                        chart.category_axis.has_title = True
                        chart.category_axis.axis_title.text_frame.text = x_label
                    except Exception as exc:  # pragma: no cover
                        logger.debug("x-axis title failed: %s", exc)
                if y_label:
                    try:
                        chart.value_axis.has_title = True
                        chart.value_axis.axis_title.text_frame.text = y_label
                    except Exception as exc:  # pragma: no cover
                        logger.debug("y-axis title failed: %s", exc)
        except Exception as exc:
            logger.warning("chart styling failed: %s", exc, exc_info=True)

    @staticmethod
    def _normalise_chart_data(
        chart_type: str, data: Any
    ) -> Tuple[List[str], List[Tuple[str, List[Any]]]]:
        """Normalise the user-supplied data dict/list into categories + series.

        Convention:
            * For ``'bar'`` / ``'line'``: the *first* series is treated as
              category labels; remaining series become chart series.
            * For ``'pie'``: the first series is categories, the second
              series is values.
        """
        if isinstance(data, dict):
            items = list(data.items())
        else:
            items = list(data)
        if not items:
            return [], []
        # First entry -> categories.
        first_name, first_values = items[0]
        categories = [str(v) for v in (first_values or [])]
        if chart_type == "pie":
            # Second series becomes the values.
            if len(items) >= 2:
                second_name, second_values = items[1]
                return categories, [(second_name or "Value", list(second_values or []))]
            return categories, [("Value", [])]
        # bar / line: subsequent items are series.
        series_list = []
        for name, values in items[1:]:
            series_list.append((str(name) if name else "Series", list(values or [])))
        return categories, series_list

    def add_image_slide(
        self,
        title: str,
        image_path: str,
        caption: Optional[str] = None,
    ) -> None:
        """Add a slide with a title, full-width image, and optional caption."""
        slide = self._new_blank_slide()
        self._add_textbox(
            slide, 0.6, 0.4, _SLIDE_W_IN - 1.2, 0.9,
            text=title, font_size=28, bold=True, color_hex=_PRIMARY_BLUE,
            align="left",
        )
        self._add_accent_bar(slide, top_in=1.25)
        if not os.path.exists(image_path):
            logger.warning("image not found: %s", image_path)
            self._add_textbox(
                slide, 0.6, 3.0, _SLIDE_W_IN - 1.2, 0.6,
                text=f"[image missing: {image_path}]",
                color_hex="D0021B",
                align="center",
            )
            return
        # Compute image dimensions: fit inside a 12 x 5.0 inch box, preserving aspect.
        try:
            from PIL import Image as PILImage  # noqa: WPS433
            with PILImage.open(image_path) as im:
                iw, ih = im.size
            aspect = ih / iw if iw else 0.75
        except Exception:
            aspect = 0.625  # 16:10 fallback.

        max_w = _SLIDE_W_IN - 1.6
        max_h = _SLIDE_H_IN - 2.8
        if caption:
            max_h -= 0.6
        # Fit inside the box.
        w_in = min(max_w, max_h / aspect) if aspect else max_w
        h_in = w_in * aspect
        left_in = (_SLIDE_W_IN - w_in) / 2.0
        top_in = 1.7 + (max_h - h_in) / 2.0
        try:
            slide.shapes.add_picture(
                image_path, self._Inches(left_in), self._Inches(top_in),
                self._Inches(w_in), self._Inches(h_in),
            )
        except Exception as exc:
            logger.error("add_picture failed: %s", exc, exc_info=True)
            return
        if caption:
            self._add_textbox(
                slide, 0.6, _SLIDE_H_IN - 0.8, _SLIDE_W_IN - 1.2, 0.6,
                text=caption, font_size=12, color_hex=_TEXT_GREY,
                align="center",
            )

    def add_table_slide(
        self,
        title: str,
        headers: Sequence[str],
        rows: Sequence[Sequence[Any]],
    ) -> None:
        """Add a slide with a styled table.

        Args:
            title: Slide title.
            headers: Column header strings.
            rows: Sequence of row sequences (each row has one cell per column).
        """
        slide = self._new_blank_slide()
        self._add_textbox(
            slide, 0.6, 0.4, _SLIDE_W_IN - 1.2, 0.9,
            text=title, font_size=28, bold=True, color_hex=_PRIMARY_BLUE,
            align="left",
        )
        self._add_accent_bar(slide, top_in=1.25)
        n_rows = len(rows) + 1
        n_cols = len(headers)
        if n_cols == 0 or n_rows == 1:
            self._add_textbox(
                slide, 0.6, 3.0, _SLIDE_W_IN - 1.2, 0.6,
                text="(no data)", color_hex=_TEXT_GREY, align="center",
            )
            return
        table_left = self._Inches(0.6)
        table_top = self._Inches(1.7)
        table_width = self._Inches(_SLIDE_W_IN - 1.2)
        table_height = self._Inches(min(_SLIDE_H_IN - 2.3, 0.35 * n_rows))
        table_shape = slide.shapes.add_table(
            n_rows, n_cols, table_left, table_top, table_width, table_height,
        )
        table = table_shape.table
        # Header row.
        for j, h in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = str(h)
            try:
                from pptx.dml.color import RGBColor  # noqa: WPS433
                cell.fill.solid()
                cell.fill.fore_color.rgb = self._hex_color(_PRIMARY_BLUE)
                para = cell.text_frame.paragraphs[0]
                for run in para.runs:
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    run.font.size = self._Pt(14)
            except Exception as exc:  # pragma: no cover
                logger.debug("header cell format failed: %s", exc)
        # Body rows.
        for i, row in enumerate(rows, start=1):
            for j, val in enumerate(row):
                if j >= n_cols:
                    break
                cell = table.cell(i, j)
                cell.text = "" if val is None else str(val)
                try:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self._hex_color(
                        "FFFFFF" if i % 2 else _BG_TINT
                    )
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.size = self._Pt(12)
                            run.font.color.rgb = self._hex_color(_TEXT_DARK)
                except Exception as exc:  # pragma: no cover
                    logger.debug("body cell format failed: %s", exc)

    def add_two_column_slide(
        self,
        title: str,
        left_title: str,
        left_items: Sequence[str],
        right_title: str,
        right_items: Sequence[str],
    ) -> None:
        """Add a two-column comparison slide.

        Args:
            title: Slide title.
            left_title: Heading above the left column.
            left_items: Bullet strings for the left column.
            right_title: Heading above the right column.
            right_items: Bullet strings for the right column.
        """
        slide = self._new_blank_slide()
        self._add_textbox(
            slide, 0.6, 0.4, _SLIDE_W_IN - 1.2, 0.9,
            text=title, font_size=28, bold=True, color_hex=_PRIMARY_BLUE,
            align="left",
        )
        self._add_accent_bar(slide, top_in=1.25)
        col_w = (_SLIDE_W_IN - 2.0) / 2.0
        left_x = 0.6
        right_x = left_x + col_w + 0.8
        top_y = 1.7
        col_h = _SLIDE_H_IN - 2.2

        # Column header bars.
        from pptx.enum.shapes import MSO_SHAPE  # noqa: WPS433
        for x, col_title in ((left_x, left_title), (right_x, right_title)):
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                self._Inches(x), self._Inches(top_y),
                self._Inches(col_w), self._Inches(0.5),
            )
            try:
                bar.fill.solid()
                bar.fill.fore_color.rgb = self._hex_color(_LIGHT_BLUE)
                bar.line.fill.background()
                bar.shadow.inherit = False
                tf = bar.text_frame
                tf.text = col_title
                for para in tf.paragraphs:
                    para.alignment = 2  # CENTER
                    for run in para.runs:
                        run.font.bold = True
                        run.font.color.rgb = self._hex_color("FFFFFF")
                        run.font.size = self._Pt(16)
            except Exception as exc:  # pragma: no cover
                logger.debug("col header bar failed: %s", exc)

        # Column bullet bodies.
        for x, items in ((left_x, left_items), (right_x, right_items)):
            tb = slide.shapes.add_textbox(
                self._Inches(x), self._Inches(top_y + 0.7),
                self._Inches(col_w), self._Inches(col_h - 0.8),
            )
            tf = tb.text_frame
            tf.word_wrap = True
            first = True
            for item in items:
                if first:
                    p = tf.paragraphs[0]
                    first = False
                else:
                    p = tf.add_paragraph()
                run = p.add_run()
                run.text = f"•  {item}"
                run.font.size = self._Pt(14)
                run.font.color.rgb = self._hex_color(_TEXT_DARK)
                try:
                    p.space_after = self._Pt(6)
                except Exception:  # pragma: no cover
                    pass

    # ------------------------------------------------------------------
    # Build
    # ------------------------------------------------------------------
    def build(self, path: str) -> str:
        """Save the deck to ``path`` and return the absolute path.

        Args:
            path: Target ``.pptx`` file path (parent dirs auto-created).
        """
        os.makedirs(os.path.dirname(os.path.abspath(path)) or ".", exist_ok=True)
        abs_path = os.path.abspath(path)
        try:
            self._prs.save(abs_path)
        except Exception as exc:
            logger.error("PPTX save failed -> %s: %s", abs_path, exc, exc_info=True)
            raise
        logger.info("PPTX built -> %s", abs_path)
        return abs_path


__all__ = ["PPTXReport"]
